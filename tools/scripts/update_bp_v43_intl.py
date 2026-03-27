"""
Surgical update of BP PPTX: database enrichment + internationalization
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import re

pptx_path = r"C:\Users\gvera\Projects\FitManager_AI_Studio\FitManager_Business_Plan_2026.pptx"
prs = Presentation(pptx_path)

TOTAL_SLIDES = len(prs.slides)
print(f"Current slides: {TOTAL_SLIDES}")

# --- Helper: add text box ---
def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=RGBColor(0x64, 0x64, 0x64), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

# --- Helper: add table row ---
def add_table_row(table, cells_text, bold_row=False, font_size=None):
    """Add a row to an existing table, copying style from the last row."""
    from pptx.oxml.ns import qn

    tbl = table._tbl
    # Clone the last row
    all_trs = tbl.findall(qn('a:tr'))
    last_tr = all_trs[-1]
    new_tr = deepcopy(last_tr)

    # Update cell text
    tcs = new_tr.findall(qn('a:tc'))
    for i, tc in enumerate(tcs):
        if i < len(cells_text):
            # Clear all paragraphs and set new text
            for p in tc.findall(qn('a:p')):
                for r in p.findall(qn('a:r')):
                    t = r.find(qn('a:t'))
                    if t is not None:
                        t.text = cells_text[i] if i == 0 or not cells_text[i].startswith('—') else cells_text[i]
                    # Set bold if needed
                    if bold_row:
                        rPr = r.find(qn('a:rPr'))
                        if rPr is not None:
                            rPr.set('b', '1')
                # Only keep first run, clear others
                runs = p.findall(qn('a:r'))
                if runs:
                    runs[0].find(qn('a:t')).text = cells_text[i]
                    for extra_run in runs[1:]:
                        p.remove(extra_run)

    tbl.append(new_tr)

# --- Helper: find text in slide ---
def find_shape_with_text(slide, search_text):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if search_text in para.text:
                    return shape, para
    return None, None

# --- Helper: insert slide after index ---
def insert_slide_after(prs, index, layout_index=6):
    """Insert a blank slide after the given index."""
    # Use blank layout
    layout = prs.slide_layouts[layout_index] if layout_index < len(prs.slide_layouts) else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)

    # Move slide to correct position
    from pptx.oxml.ns import qn as _qn
    sldIdLst = prs.part._element.find(_qn('p:sldIdLst'))
    sldIds = list(sldIdLst)
    # The new slide is at the end, move it to index+1
    new_sldId = sldIds[-1]
    sldIdLst.remove(new_sldId)
    # Insert after the slide at `index`
    ref_sldId = sldIds[index]
    ref_idx = list(sldIdLst).index(ref_sldId)
    sldIdLst.insert(ref_idx + 1, new_sldId)

    return slide

# ============================================================
# 1. Slide 8 (index 7): Add database enrichment note at bottom
# ============================================================
slide8 = prs.slides[7]
add_textbox(slide8, 0.5, 4.85, 9.0, 0.5,
    "Il database è progettato per crescere con il contributo di professionisti del settore. "
    "Esercizi validati da migliaia di ore di pratica 1:1, varianti testate su clienti reali con patologie specifiche.",
    font_size=10, color=RGBColor(0x80, 0x80, 0x80))
print("1. Slide 8: database enrichment note added")

# ============================================================
# 2. Slide 12 (index 11): Add database asset text
#    Actually slide 12 is "Confronto con le alternative" which has the table
#    The "Quanto costa" is slide 13. Let me add to slide 12 since it mentions
#    "Nessun prodotto..." at the bottom
# ============================================================
slide12 = prs.slides[11]
add_textbox(slide12, 0.5, 4.55, 9.0, 0.55,
    "Il database esercizi non è statico — arricchito da professionisti certificati con varianti "
    "e progressioni validate sul campo. Un asset che cresce in valore nel tempo.",
    font_size=10, color=RGBColor(0x80, 0x80, 0x80))
print("2. Slide 12: database asset note added")

# ============================================================
# 3. Insert NEW slide after slide 15 (index 14): Mercato internazionale
# ============================================================
new_slide = insert_slide_after(prs, 14)

# Set background
bg = new_slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xF5, 0xF6, 0xFA)

# Section label
add_textbox(new_slide, 0.5, 0.3, 4.0, 0.3,
    "SEZIONE 5 — IL MERCATO", font_size=10, bold=True,
    color=RGBColor(0x3B, 0x82, 0xF6))

# Title
add_textbox(new_slide, 0.5, 0.6, 9.0, 0.5,
    "Oltre l'Italia: il mercato internazionale",
    font_size=28, bold=True, color=RGBColor(0x1A, 0x1F, 0x36))

# Intro text
add_textbox(new_slide, 0.5, 1.15, 9.0, 0.4,
    "L'architettura è progettata per l'internazionalizzazione. "
    "L'approccio è incrementale, per blocchi di complessità crescente.",
    font_size=12, color=RGBColor(0x4B, 0x55, 0x63))

# 4 blocks
blocks = [
    ("Blocco 1 — Interfaccia e core", "BASSA", "Traduzione UI, template, formati data/valuta. Poche settimane.", 1.65),
    ("Blocco 2 — Esercizi e Safety Engine", "MEDIA", "Scienza universale (biomeccanica, fisiologia). Adattamento terminologico professionale.", 2.35),
    ("Blocco 3 — Nutrizione e compliance", "ALTA", "Database locali (USDA, McCance & Widdowson, EuroFIR). Anno 2+.", 3.05),
    ("Blocco 4 — Moduli fiscali", "VARIABILE", "Pagamenti, fatturazione per giurisdizione. Configurazione modulare.", 3.75),
]

for title, complexity, desc, y in blocks:
    # Accent bar
    new_slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(y), Inches(0.06), Inches(0.55),
        )
    accent = new_slide.shapes[-1]
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    accent.line.fill.background()

    # Title + complexity
    add_textbox(new_slide, 0.7, y, 5.0, 0.25,
        title, font_size=13, bold=True, color=RGBColor(0x1E, 0x29, 0x3B))
    add_textbox(new_slide, 5.8, y, 1.2, 0.25,
        complexity, font_size=10, bold=True,
        color=RGBColor(0xF5, 0x9E, 0x0B) if complexity in ("MEDIA", "ALTA") else RGBColor(0x10, 0xB9, 0x81),
        alignment=PP_ALIGN.RIGHT)
    # Description
    add_textbox(new_slide, 0.7, y + 0.25, 8.3, 0.28,
        desc, font_size=11, color=RGBColor(0x64, 0x74, 0x8B))

# Footer
add_textbox(new_slide, 0.5, 4.55, 9.0, 0.5,
    "L'italiano resta la priorità Anno 1. L'inglese è la leva di crescita "
    "che trasforma un prodotto verticale italiano in un prodotto scalabile.",
    font_size=11, color=RGBColor(0x64, 0x74, 0x8B))

# Page number
add_textbox(new_slide, 8.5, 5.15, 1.0, 0.3,
    "16 / 56", font_size=9, color=RGBColor(0x9C, 0xA3, 0xAF), alignment=PP_ALIGN.RIGHT)

print("3. New slide inserted: Mercato internazionale (after slide 15)")

# ============================================================
# 4. Slide ~39 (Team, Partner): Add database contribution
#    After insertion, this is now index 39 (was 38)
# ============================================================
# Find slide with "INDUSTRY PARTNER" and partner role text
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "Valida posizionamento PT Evoluto" in full_text or "valida il posizionamento" in full_text.lower():
                # Add to this shape
                for para in shape.text_frame.paragraphs:
                    if "PT Evoluto" in para.text and "posizionamento" in para.text.lower():
                        # Append to this paragraph
                        original = para.text
                        if "database" not in original.lower():
                            for run in para.runs:
                                if "PT Evoluto" in run.text:
                                    run.text = run.text.rstrip()
                                    if not run.text.endswith('.'):
                                        run.text += "."
                            # Add new paragraph
                            new_p = shape.text_frame.add_paragraph()
                            new_p.text = "Contribuisce all'arricchimento del database esercizi con varianti e progressioni validate dalla propria esperienza pratica."
                            new_p.font.size = para.runs[0].font.size if para.runs else Pt(12)
                            new_p.font.color.rgb = para.runs[0].font.color.rgb if para.runs and para.runs[0].font.color.rgb else RGBColor(0x64, 0x64, 0x64)
                            print(f"4. Slide {idx+1}: Partner database contribution added")
                            break
                break

# ============================================================
# 5. Slide ~42 (Piano di crescita): Add internationalization to Fase 4
#    After insertion, index shifted by 1
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "Certificazione PT Evoluto" in full_text and "Fase 4" in full_text:
                for para in shape.text_frame.paragraphs:
                    if "Certificazione PT Evoluto" in para.text:
                        for run in para.runs:
                            if "Certificazione" in run.text:
                                run.text = run.text.rstrip('. ')
                                run.text += ". Lancio versione inglese (Blocchi 1-2). Primi clienti internazionali."
                                print(f"5. Slide {idx+1}: Internationalization added to Fase 4")
                                break
                        break
                break

# ============================================================
# 6. Slide ~44 (A1 table): Update Allenamento row
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for row_idx in range(len(table.rows)):
                cell_text = table.cell(row_idx, 0).text.strip()
                if cell_text == "Allenamento":
                    cell = table.cell(row_idx, 1)
                    current = cell.text
                    if "arricchimento" not in current:
                        # Append to existing text
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if "export PDF" in run.text:
                                    run.text += " Database progettato per arricchimento continuo con professionisti certificati."
                                    print(f"6. Slide {idx+1}: A1 Allenamento row updated")
                                    break

# ============================================================
# 7. Slide ~46 (A4 table): Add DB1 and INT1 rows
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            # Check if this is the A4 table (has "Cod." or "P4" in first column)
            has_p4 = False
            for row_idx in range(len(table.rows)):
                if table.cell(row_idx, 0).text.strip() in ("P4", "Cod."):
                    has_p4 = True
                    break
            if has_p4:
                # Check if already has DB1
                has_db1 = False
                for row_idx in range(len(table.rows)):
                    if table.cell(row_idx, 0).text.strip() == "DB1":
                        has_db1 = True
                        break
                if not has_db1:
                    add_table_row(table, ["DB1", "Professionisti disponibili a contribuire al database", "Ipotesi forte", "Primo contributo partnership/POC"])
                    add_table_row(table, ["INT1", "Versione inglese core (Blocchi 1-2) entro 6 mesi", "Ipotesi", "Valutazione tecnica in corso"])
                    print(f"7. Slide {idx+1}: A4 table - DB1 and INT1 rows added")

# ============================================================
# 8. Update page numbers (55 -> 56)
# ============================================================
count = 0
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                # Match pattern "N / 55"
                match = re.match(r'^(\d+)\s*/\s*55$', text)
                if match:
                    old_num = int(match.group(1))
                    # After slide 15, shift by 1
                    new_num = old_num if old_num <= 15 else old_num + 1
                    for run in para.runs:
                        if re.match(r'^\d+\s*/\s*55$', run.text.strip()):
                            run.text = f"{new_num} / 56"
                            count += 1
                            break
print(f"8. Page numbers updated: {count} slides")

# Save
prs.save(pptx_path)
print(f"\nSaved to {pptx_path}")
print(f"Total slides: {len(prs.slides)}")
