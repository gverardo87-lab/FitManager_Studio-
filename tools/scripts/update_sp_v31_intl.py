"""
Surgical update of Strategy Plan PPTX: internationalization + database enrichment
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
import re

pptx_path = r"C:\Users\gvera\Projects\FitManager_AI_Studio\FitManager_Strategy_Plan_2026.pptx"
prs = Presentation(pptx_path)
print(f"Current slides: {len(prs.slides)}")

# --- Helpers ---
def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=RGBColor(0x64, 0x64, 0x64), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def insert_slide_after(prs, index, layout_index=6):
    layout = prs.slide_layouts[layout_index] if layout_index < len(prs.slide_layouts) else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
    sldIds = list(sldIdLst)
    new_sldId = sldIds[-1]
    sldIdLst.remove(new_sldId)
    ref_sldId = sldIds[index]
    ref_idx = list(sldIdLst).index(ref_sldId)
    sldIdLst.insert(ref_idx + 1, new_sldId)
    return slide

def setup_light_slide(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF6, 0xFA)

def add_accent_bar(slide, x, y, h, color=RGBColor(0x3B, 0x82, 0xF6)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_table_row(table, cells_text):
    tbl = table._tbl
    all_trs = tbl.findall(qn('a:tr'))
    last_tr = all_trs[-1]
    new_tr = deepcopy(last_tr)
    tcs = new_tr.findall(qn('a:tc'))
    for i, tc in enumerate(tcs):
        if i < len(cells_text):
            for p in tc.findall(qn('a:p')):
                runs = p.findall(qn('a:r'))
                if runs:
                    runs[0].find(qn('a:t')).text = cells_text[i]
                    for extra in runs[1:]:
                        p.remove(extra)
    tbl.append(new_tr)

# ============================================================
# 1. Insert NEW slide after slide 6 (index 5): "La direzione internazionale"
# ============================================================
s1 = insert_slide_after(prs, 5)
setup_light_slide(s1)
add_textbox(s1, 0.5, 0.3, 5.0, 0.3, "SEZIONE 1 — LE 5 EVOLUZIONI", font_size=10, bold=True, color=RGBColor(0x3B, 0x82, 0xF6))
add_textbox(s1, 0.5, 0.65, 9.0, 0.5, "La direzione internazionale", font_size=28, bold=True, color=RGBColor(0x1A, 0x1F, 0x36))
add_accent_bar(s1, 0.5, 1.3, 1.8)
add_textbox(s1, 0.7, 1.3, 8.5, 1.8,
    "FitManager nasce per il mercato italiano, ma l'architettura è pronta per l'internazionalizzazione. "
    "La versione inglese è in roadmap: interfaccia, database esercizi e Safety Engine entro i primi 6 mesi "
    "(la scienza dell'allenamento è universale); nutrizione e compliance locale per i mercati target "
    "nella fase di espansione.",
    font_size=14, color=RGBColor(0x37, 0x41, 0x51))
add_textbox(s1, 0.5, 3.5, 9.0, 0.8,
    "L'italiano resta la priorità dell'Anno 1 — l'inglese è la leva che trasforma la crescita "
    "da lineare a esponenziale.",
    font_size=13, bold=True, color=RGBColor(0x3B, 0x82, 0xF6))
print("1. New slide inserted after 6: La direzione internazionale")

# ============================================================
# 2. Insert NEW slide after slide ~11 (mental coaching): "PT Evoluto non ha confini"
#    After insertion above, mental coaching is now at index ~12 (was 11)
#    Let's find it by content
# ============================================================
mental_idx = None
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "mental coaching" in full.lower() and "collegamento" in full.lower():
                mental_idx = idx
                break
    if mental_idx is not None:
        break

if mental_idx is not None:
    s2 = insert_slide_after(prs, mental_idx)
    setup_light_slide(s2)
    add_textbox(s2, 0.5, 0.3, 4.0, 0.3, "SEZIONE 2", font_size=10, bold=True, color=RGBColor(0x3B, 0x82, 0xF6))
    add_textbox(s2, 0.5, 0.65, 9.0, 0.5, "Il PT Evoluto non ha confini geografici", font_size=28, bold=True, color=RGBColor(0x1A, 0x1F, 0x36))
    add_accent_bar(s2, 0.5, 1.3, 2.0)
    add_textbox(s2, 0.7, 1.3, 8.5, 2.0,
        "Il concetto di Personal Trainer Evoluto non è italiano — è universale. "
        "Un PT a Göteborg ha gli stessi problemi di uno a Genova: gestione clienti artigianale, "
        "nessuno strumento scientifico integrato, rischio errori clinici. "
        "La differenza è che in Scandinavia e nel mercato anglofono la propensione a pagare "
        "per strumenti professionali è più alta e il mercato è più maturo. "
        "La versione inglese di FitManager non è un adattamento — è un'espansione naturale "
        "di un prodotto che risolve un problema universale.",
        font_size=14, color=RGBColor(0x37, 0x41, 0x51))
    print(f"2. New slide inserted after {mental_idx+1}: PT Evoluto senza confini")
else:
    print("2. WARNING: Could not find mental coaching slide")

# ============================================================
# 3. Insert NEW slide before "Il volano completo": "L'apertura internazionale"
#    Find the slide with "Il volano completo" title
# ============================================================
volano_idx = None
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "volano completo" in full.lower() and "circolare" in full.lower():
                volano_idx = idx
                break
    if volano_idx is not None:
        break

if volano_idx is not None:
    s3 = insert_slide_after(prs, volano_idx - 1)  # Insert before volano
    setup_light_slide(s3)
    add_textbox(s3, 0.5, 0.3, 4.0, 0.3, "SEZIONE 4", font_size=10, bold=True, color=RGBColor(0x3B, 0x82, 0xF6))
    add_textbox(s3, 0.5, 0.65, 9.0, 0.5, "L'apertura internazionale", font_size=28, bold=True, color=RGBColor(0x1A, 0x1F, 0x36))

    timeline = [
        ("Mesi 1-3", "Focus totale POC italiana. Nessun lavoro sull'internazionalizzazione.", 1.3),
        ("Mesi 4-6", "Traduzione interfaccia e database esercizi (Blocchi 1-2). Primi contatti esplorativi mercato anglofono.", 2.0),
        ("Mesi 7-12", "Versione inglese core disponibile. 5-10 utenti pilota internazionali selezionati.", 2.7),
        ("Anno 2+", "Lancio strutturato. Database nutrizionale localizzato. Pricing specifico per mercato.", 3.4),
    ]
    for label, desc, y in timeline:
        add_accent_bar(s3, 0.5, y, 0.55)
        add_textbox(s3, 0.7, y, 1.5, 0.25, label, font_size=14, bold=True, color=RGBColor(0x1E, 0x29, 0x3B))
        add_textbox(s3, 2.3, y, 7.0, 0.55, desc, font_size=12, color=RGBColor(0x4B, 0x55, 0x63))

    add_textbox(s3, 0.5, 4.3, 9.0, 0.5,
        "Ogni feature sviluppata per il mercato italiano viene sviluppata bilingue da quel momento in poi.",
        font_size=11, color=RGBColor(0x64, 0x74, 0x8B))
    print(f"3. New slide inserted before {volano_idx+1}: L'apertura internazionale")
else:
    print("3. WARNING: Could not find volano completo slide")

# ============================================================
# 4. Find Scenario 3 slide: Add international leverage text
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "masterclass diventano eventi" in full.lower():
                for para in shape.text_frame.paragraphs:
                    if "masterclass diventano eventi" in para.text.lower():
                        for run in para.runs:
                            if "masterclass diventano" in run.text.lower() and "internazional" not in run.text.lower():
                                run.text = run.text.rstrip('. ') + ". Se il partner opera in contesto internazionale, lo Scenario 3 include l'apertura di mercati esteri. La versione inglese (mesi 7-12) diventa la leva per network internazionali."
                                print(f"4. Slide {idx+1}: Scenario 3 international text added")
                                break
                        break
                break

# ============================================================
# 5. Find Struttura compenso table: Add "Contributo al database" row
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            # Check if this is the compenso table (has "25%" or "Ricavi POC")
            is_compenso = False
            for row_idx in range(len(table.rows)):
                try:
                    if table.cell(row_idx, 0).text.strip() == "Ricavi POC":
                        is_compenso = True
                        break
                except:
                    pass
            if is_compenso:
                # Check if already has "Contributo"
                has_it = False
                for row_idx in range(len(table.rows)):
                    if "Contributo" in table.cell(row_idx, 0).text:
                        has_it = True
                        break
                if not has_it:
                    add_table_row(table, ["Contributo al database", "Da definire", "Esercizi, varianti, progressioni validate", "Arricchisce l'asset core"])
                    print(f"5. Slide {idx+1}: Compenso table - Contributo al database row added")
                break

# ============================================================
# 6. Find Marketing mesi 7-12 slide: Add "Primi contenuti in inglese"
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "podcast tour" in full.lower() and "comment marketing" in full.lower():
                # Add text box at bottom
                add_textbox(slide, 0.5, 4.5, 9.0, 0.6,
                    "Primi contenuti in inglese. Quando la versione inglese è disponibile: "
                    "1 post LinkedIn/settimana in inglese, landing page dedicata, 2-3 demo in inglese. "
                    "Investimento minimo ma posiziona FitManager come prodotto internazionale.",
                    font_size=11, color=RGBColor(0x4B, 0x55, 0x63))
                print(f"6. Slide {idx+1}: Marketing - contenuti inglese added")
                break

# ============================================================
# 7. Find Rischi table: Add internationalization risk row
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            # Check if this is the risk table (has "Rischio" header or "trainer non percepiscono")
            is_risk = False
            for row_idx in range(len(table.rows)):
                try:
                    cell_text = table.cell(row_idx, 0).text.strip().lower()
                    if "trainer non percepiscono" in cell_text or "rischio" == cell_text:
                        is_risk = True
                        break
                except:
                    pass
            if is_risk:
                has_intl = False
                for row_idx in range(len(table.rows)):
                    if "internazionalizzazione" in table.cell(row_idx, 0).text.lower():
                        has_intl = True
                        break
                if not has_intl:
                    add_table_row(table, [
                        "Internazionalizzazione prematura",
                        "Blocchi 1-2 costano settimane, non mesi. Se mercato non risponde, costo contenuto. Blocco 3 solo con domanda validata.",
                        "Mesi 7-12"
                    ])
                    print(f"7. Slide {idx+1}: Risk table - internationalization row added")
                break

# ============================================================
# 8. Update page numbers (47 -> 50)
# ============================================================
new_total = len(prs.slides)
count = 0
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                match = re.match(r'^(\d+)\s*/\s*47$', text)
                if match:
                    old_num = int(match.group(1))
                    # Calculate new number based on insertions
                    # Inserted after 6, after ~12, before volano (~22)
                    # Shift logic: slides after each insertion point get +1
                    new_num = old_num
                    if old_num > 6:
                        new_num += 1
                    if old_num > 11:  # mental coaching was ~11
                        new_num += 1
                    if old_num > 20:  # before volano was ~21
                        new_num += 1
                    for run in para.runs:
                        if re.match(r'^\d+\s*/\s*47$', run.text.strip()):
                            run.text = f"{new_num} / {new_total}"
                            count += 1
                            break
print(f"8. Page numbers updated: {count} slides (total now {new_total})")

# Save
prs.save(pptx_path)
print(f"\nSaved to {pptx_path}")
