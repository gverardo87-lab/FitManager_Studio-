"""
Update FitManager Business Plan PPTX from v4.2 to v4.3
Adds WhatsApp/Email communication integration content.
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
import os

PPTX_PATH = os.path.join(os.path.dirname(__file__), '..', '..', 'FitManager_Business_Plan_2026.pptx')
PPTX_PATH = os.path.abspath(PPTX_PATH)

prs = Presentation(PPTX_PATH)

# ── Helpers ──

def set_cell_fill(cell, color_hex):
    """Set solid fill on a table cell."""
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = tc.makeelement(qn('a:tcPr'), {})
        tc.insert(0, tcPr)
    # Remove existing fill
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': color_hex})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def set_cell_text(cell, text, font_size, font_color, bold=None):
    """Set text in a table cell with formatting."""
    cell.text = ""
    para = cell.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    if bold is not None:
        run.font.bold = bold


def add_table_row(table, cells_data, row_fill_color, is_header=False):
    """Add a row to a table by cloning the last row XML and modifying it."""
    # Clone last row
    tbl = table._tbl
    last_tr = tbl.findall(qn('a:tr'))[-1]
    new_tr = deepcopy(last_tr)
    tbl.append(new_tr)

    # Get the new row index
    new_row_idx = len(table.rows) - 1
    row = table.rows[new_row_idx]
    h_val = last_tr.get('h')
    if h_val is not None:
        row.height = int(h_val)

    # Set cell content and fill
    for ci, (text, font_size, font_color, bold) in enumerate(cells_data):
        if ci < len(table.columns):
            cell = table.cell(new_row_idx, ci)
            set_cell_text(cell, text, font_size, font_color, bold)
            set_cell_fill(cell, row_fill_color)

    return new_row_idx


def update_page_numbers(prs, from_slide_idx, total_new):
    """Update page number text boxes (format 'N / 54' -> 'N / 55')."""
    new_total = 54 + total_new
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if '/' in text and text.replace(' ', '').replace('/', '').isdigit():
                        parts = text.split('/')
                        if len(parts) == 2:
                            try:
                                page_num = int(parts[0].strip())
                                old_total = int(parts[1].strip())
                                if old_total == 54:
                                    # Adjust page number for slides after insertion
                                    if idx >= from_slide_idx:
                                        page_num += total_new
                                    new_text = f"{page_num} / {new_total}"
                                    if para.runs:
                                        para.runs[0].text = new_text
                                        # Clear other runs if any
                                        for r in para.runs[1:]:
                                            r.text = ""
                            except ValueError:
                                pass


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title: v4.2 → v4.3, 26 marzo → 27 marzo
# ═══════════════════════════════════════════════════════════════════
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'v4.2' in run.text or '4.2' in run.text:
                    run.text = run.text.replace('4.2', '4.3')
                if '26 marzo' in run.text:
                    run.text = run.text.replace('26 marzo', '27 marzo')

print("[OK] Slide 1 updated (v4.3, 27 marzo)")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — IL COMPETITOR VERO: add WhatsApp paragraph
# ═══════════════════════════════════════════════════════════════════
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.name == 'Text 4':
        tf = shape.text_frame
        # Add blank line then new paragraph
        para_new = tf.add_paragraph()
        para_new.space_before = Pt(12)
        run = para_new.add_run()
        run.text = (
            "WhatsApp \u00e8 il centro del workflow di ogni trainer. \u00c8 l\u00ec che conferma "
            "appuntamenti, manda schede, sollecita pagamenti, tiene i rapporti con i clienti. "
            "Ma lo fa manualmente \u2014 riscrive gli stessi messaggi 20 volte al giorno, dimentica "
            "promemoria, perde il filo delle conversazioni. Qualunque soluzione che chieda al "
            "trainer di abbandonare WhatsApp parte sconfitta. La soluzione giusta non sostituisce "
            "WhatsApp \u2014 lo potenzia."
        )
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        run.font.bold = None

        # Resize the text box to accommodate extra text
        shape.height = shape.height + Emu(600000)

print("[OK] Slide 6 updated (WhatsApp paragraph)")


# ═══════════════════════════════════════════════════════════════════
# NEW SLIDE — Insert after slide 8 (before "La FitManager Box")
# ═══════════════════════════════════════════════════════════════════

# Create new slide using the only layout available
layout = prs.slide_layouts[0]

# We need to insert at position 8 (after current slide 8, before slide 9)
# python-pptx doesn't have an insert method, so we add and then move
new_slide = prs.slides.add_slide(layout)

# Set dark background matching section slides (0F1923)
bg = new_slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x0F, 0x19, 0x23)

# Remove any default placeholder shapes
for sp in list(new_slide.shapes):
    sp_elem = sp._element
    sp_elem.getparent().remove(sp_elem)

# Section header
txBox = new_slide.shapes.add_textbox(Emu(457200), Emu(228600), Emu(3657600), Emu(274320))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "SEZIONE 3 \u2014 IL SOFTWARE"
run.font.size = Pt(8)
run.font.color.rgb = RGBColor(0x14, 0xB8, 0xA6)
run.font.bold = True
txBox.fill.background()

# Title
txBox2 = new_slide.shapes.add_textbox(Emu(457200), Emu(502920), Emu(8229600), Emu(502920))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = "Comunicazione integrata con i clienti"
run2.font.size = Pt(28)
run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
run2.font.bold = True
txBox2.fill.background()

# Divider line
div = new_slide.shapes.add_shape(
    1, Emu(457200), Emu(1051560), Emu(1097280), Emu(36576)
)
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0x14, 0xB8, 0xA6)
div.line.fill.background()

# Block 1: WhatsApp semi-automatico
b1_title = new_slide.shapes.add_textbox(Emu(457200), Emu(1280160), Emu(3840480), Emu(274320))
tf = b1_title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "WhatsApp semi-automatico."
run.font.size = Pt(16)
run.font.color.rgb = RGBColor(0x14, 0xB8, 0xA6)
run.font.bold = True
b1_title.fill.background()

b1_desc = new_slide.shapes.add_textbox(Emu(457200), Emu(1554480), Emu(3840480), Emu(1005840))
tf = b1_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = (
    "15 template pre-compilati per ogni situazione: sollecito rata, conferma appuntamento, "
    "scheda pronta, benvenuto, auguri compleanno, rinnovo. Un click apre WhatsApp con il messaggio "
    "gi\u00e0 pronto \u2014 il trainer rivede e invia. Nessun bot, nessun automatismo che il cliente "
    "percepisce come impersonale. Fire-and-forget: il click logga la comunicazione nel CRM."
)
run.font.size = Pt(13)
run.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
b1_desc.fill.background()

# Block 2: Email automatiche
b2_title = new_slide.shapes.add_textbox(Emu(4846320), Emu(1280160), Emu(3840480), Emu(274320))
tf = b2_title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Email automatiche."
run.font.size = Pt(16)
run.font.color.rgb = RGBColor(0x14, 0xB8, 0xA6)
run.font.bold = True
b2_title.fill.background()

b2_desc = new_slide.shapes.add_textbox(Emu(4846320), Emu(1554480), Emu(3840480), Emu(1005840))
tf = b2_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = (
    "SMTP integrato per comunicazioni formali: conferma contratto, riepilogo pagamenti, "
    "invio scheda PDF. Il trainer configura il proprio indirizzo email una volta. "
    "Il sistema genera e invia automaticamente. Ogni email viene loggata nel registro "
    "comunicazioni del cliente."
)
run.font.size = Pt(13)
run.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
b2_desc.fill.background()

# Centro Comunicazioni block
cc_title = new_slide.shapes.add_textbox(Emu(457200), Emu(2743200), Emu(8229600), Emu(274320))
tf = cc_title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Centro Comunicazioni"
run.font.size = Pt(16)
run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
run.font.bold = True
cc_title.fill.background()

cc_desc = new_slide.shapes.add_textbox(Emu(457200), Emu(3017520), Emu(8229600), Emu(548640))
tf = cc_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = (
    "Pagina dedicata con 2 tab: Invia (seleziona clienti, scegli template, invio multiplo "
    "sequenziale) e Registro (timeline di tutte le comunicazioni, filtri per cliente e periodo). "
    "Alert compleanni automatici in dashboard. Invio multiplo con stepper visuale."
)
run.font.size = Pt(13)
run.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
cc_desc.fill.background()

# Footer quote
footer_q = new_slide.shapes.add_textbox(Emu(457200), Emu(3931920), Emu(8229600), Emu(548640))
tf = footer_q.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = (
    "Il trainer non deve cambiare il modo in cui comunica con i clienti. "
    "FitManager si inserisce nel suo workflow esistente e lo rende professionale."
)
run.font.size = Pt(14)
run.font.color.rgb = RGBColor(0xF9, 0x61, 0x67)
run.font.bold = True
footer_q.fill.background()

# Page number placeholder (will be updated later)
pn = new_slide.shapes.add_textbox(Emu(7772400), Emu(4754880), Emu(1097280), Emu(274320))
tf = pn.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
run = p.add_run()
run.text = "9 / 55"
run.font.size = Pt(9)
run.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
pn.fill.background()

# Move the new slide from the end (position 54) to position 8 (after slide 8)
# python-pptx: slides are in sldIdLst, we need to move the last entry
sldIdLst = prs._element.find(qn('p:sldIdLst'))
sldId_entries = list(sldIdLst)
new_entry = sldId_entries[-1]  # The one we just added
sldIdLst.remove(new_entry)
# Insert after index 8 (0-based), which is after slide 8
ref_entry = sldId_entries[8]  # This is slide 9 (0-based index 8)
ref_entry.addprevious(new_entry)

print("[OK] New slide 9 inserted (Comunicazione integrata)")


# ═══════════════════════════════════════════════════════════════════
# Update page numbers BEFORE modifying other slides
# (because slide indices shifted after insertion)
# ═══════════════════════════════════════════════════════════════════
update_page_numbers(prs, from_slide_idx=8, total_new=1)
print("[OK] Page numbers updated (X / 55)")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 (was 11) — Comparison table: add row
# ═══════════════════════════════════════════════════════════════════
slide12 = prs.slides[11]  # was index 10, now 11 after insertion
for shape in slide12.shapes:
    if shape.has_table:
        table = shape.table
        # Find "Accesso mobile/tablet" row index (should be row 2)
        insert_after = None
        for ri, row in enumerate(table.rows):
            if 'Accesso mobile' in table.cell(ri, 0).text:
                insert_after = ri
                break

        # Determine fill color: alternating F1F5F9 / FFFFFF
        # Current rows: 0=header, 1=data, 2=data, 3=data, 4=data, 5=data, 6=data
        # New row will be row 7, so it should follow the alternation pattern
        # Rows: 1=F1F5F9, 2=FFFFFF, 3=F1F5F9, 4=FFFFFF, 5=F1F5F9, 6=FFFFFF
        # Row 7 = F1F5F9

        cells_data = [
            ("Comunicazione cliente integrata", Pt(10), RGBColor(0x1E, 0x29, 0x3B), True),
            ("S\u00ec (WhatsApp + email)", Pt(10), RGBColor(0x0D, 0x94, 0x88), True),
            ("Parziale (notifiche in-app)", Pt(10), RGBColor(0x1E, 0x29, 0x3B), None),
            ("Parziale (notifiche in-app)", Pt(10), RGBColor(0x1E, 0x29, 0x3B), None),
            ("Manuale", Pt(10), RGBColor(0x1E, 0x29, 0x3B), None),
        ]
        add_table_row(table, cells_data, "F1F5F9")
        print("[OK] Slide 12 (was 11) table updated (new row)")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 20 (was 19) — Il funnel: add text
# ═══════════════════════════════════════════════════════════════════
slide20 = prs.slides[19]  # was index 18, now 19
for shape in slide20.shapes:
    if shape.name == 'Text 13':
        tf = shape.text_frame
        para_new = tf.add_paragraph()
        para_new.space_before = Pt(10)
        run = para_new.add_run()
        run.text = (
            "La demo include un momento di dimostrazione live: il trainer vede il proprio nome, "
            "il proprio cliente e un messaggio WhatsApp pronto da inviare con un click. "
            "Questo \u00e8 il punto di conversione emotiva."
        )
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xF9, 0x61, 0x67)
        run.font.bold = True
        # Expand shape to fit
        shape.height = shape.height + Emu(400000)
        break

print("[OK] Slide 20 (was 19) updated (funnel WhatsApp moment)")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 24 (was 23) — POC intro: add text
# ═══════════════════════════════════════════════════════════════════
slide24 = prs.slides[23]  # was index 22, now 23
for shape in slide24.shapes:
    if shape.name == 'Text 4':
        tf = shape.text_frame
        para_new = tf.add_paragraph()
        para_new.space_before = Pt(8)
        run = para_new.add_run()
        run.text = (
            "La comunicazione WhatsApp \u00e8 attiva dal giorno uno. \u00c8 la prima feature che i "
            "Fondatori usano \u2014 non serve imparare il sistema per trarne valore immediato."
        )
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0x14, 0xB8, 0xA6)
        run.font.bold = True
        shape.height = shape.height + Emu(350000)
        break

print("[OK] Slide 24 (was 23) updated (POC WhatsApp)")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 26 (was 25) — Protocollo Fase A: add KPI text
# ═══════════════════════════════════════════════════════════════════
slide26 = prs.slides[25]  # was index 24, now 25
for shape in slide26.shapes:
    if shape.name == 'Text 7':  # Fase A description
        tf = shape.text_frame
        para_new = tf.add_paragraph()
        para_new.space_before = Pt(4)
        run = para_new.add_run()
        run.text = (
            "KPI di attivazione: primo messaggio WhatsApp pre-compilato inviato "
            "entro 24 ore dall'installazione."
        )
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0x14, 0xB8, 0xA6)
        run.font.bold = True
        shape.height = shape.height + Emu(200000)
        break

print("[OK] Slide 26 (was 25) updated (Fase A KPI)")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 27 (was 26) — Metriche POC table: add 2 rows
# ═══════════════════════════════════════════════════════════════════
slide27 = prs.slides[26]  # was index 25, now 26
for shape in slide27.shapes:
    if shape.has_table:
        table = shape.table
        # Row 9 (odd → F1F5F9), Row 10 (even → FFFFFF)
        cells1 = [
            ("Messaggi WA pre-compilati/settimana", Pt(9), RGBColor(0x1E, 0x29, 0x3B), None),
            ("\u2014", Pt(9), RGBColor(0x1E, 0x29, 0x3B), None),
            ("15+", Pt(9), RGBColor(0x1E, 0x29, 0x3B), None),
            ("Adozione comunicazione", Pt(9), RGBColor(0x1E, 0x29, 0x3B), None),
        ]
        add_table_row(table, cells1, "F1F5F9")

        cells2 = [
            ("Tempo risparmiato comunicazione", Pt(9), RGBColor(0x1E, 0x29, 0x3B), None),
            ("\u2014", Pt(9), RGBColor(0x1E, 0x29, 0x3B), None),
            ("30+ min/giorno", Pt(9), RGBColor(0x1E, 0x29, 0x3B), None),
            ("Valore feature", Pt(9), RGBColor(0x1E, 0x29, 0x3B), None),
        ]
        add_table_row(table, cells2, "FFFFFF")
        break

print("[OK] Slide 27 (was 26) table updated (2 new metric rows)")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 42 (was 41) — Piano crescita, Fase 2 description
# ═══════════════════════════════════════════════════════════════════
slide42 = prs.slides[41]  # was index 40, now 41
for shape in slide42.shapes:
    if shape.name == 'Text 14':  # Fase 2 description
        tf = shape.text_frame
        # Append to existing paragraph or add new
        para_new = tf.add_paragraph()
        run = para_new.add_run()
        run.text = (
            "Il pitch parte dal tangibile: \u201cun click e il promemoria parte su WhatsApp. "
            "Vuoi vedere cosa fa il resto?\u201d"
        )
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x14, 0xB8, 0xA6)
        run.font.bold = True
        break

print("[OK] Slide 42 (was 41) updated (Fase 2 WhatsApp pitch)")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 44 (was 43) — A1 table: add row
# ═══════════════════════════════════════════════════════════════════
slide44 = prs.slides[43]  # was index 42, now 43
for shape in slide44.shapes:
    if shape.has_table:
        table = shape.table
        # Row 7 (even index in data = FFFFFF)
        cells = [
            ("Comunicazione", Pt(10), RGBColor(0x1E, 0x29, 0x3B), True),
            ("WhatsApp semi-auto (wa.me), email SMTP auto, template", Pt(10), RGBColor(0x1E, 0x29, 0x3B), None),
            ("Completo", Pt(10), RGBColor(0x0D, 0x94, 0x88), True),
        ]
        # Alternation: rows 1=F1F5F9(?), let's check: 7 data rows, new is 8th
        # Current: 0=header, 1-6=data. Row 7 = odd data → F1F5F9
        add_table_row(table, cells, "F1F5F9")
        break

print("[OK] Slide 44 (was 43) A1 table updated")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 46 (was 45) — A4 table: add 2 rows
# ═══════════════════════════════════════════════════════════════════
slide46 = prs.slides[45]  # was index 44, now 45
for shape in slide46.shapes:
    if shape.has_table:
        table = shape.table
        # Current: 10 rows (0=header + 9 data). New rows 10, 11.
        # Row 10 = F1F5F9 (even data index), Row 11 = FFFFFF
        cells1 = [
            ("WA1", Pt(8.5), RGBColor(0x1E, 0x29, 0x3B), True),
            ("WhatsApp semi-auto riduce frizione iniziale", Pt(8.5), RGBColor(0x1E, 0x29, 0x3B), None),
            ("Ipotesi forte", Pt(8.5), RGBColor(0x1E, 0x29, 0x3B), None),
            ("KPI attivazione POC", Pt(8.5), RGBColor(0x1E, 0x29, 0x3B), None),
        ]
        add_table_row(table, cells1, "F1F5F9")

        cells2 = [
            ("WA2", Pt(8.5), RGBColor(0x1E, 0x29, 0x3B), True),
            ("Trainer usano 15+ msg pre-compilati/sett", Pt(8.5), RGBColor(0x1E, 0x29, 0x3B), None),
            ("Ipotesi", Pt(8.5), RGBColor(0x1E, 0x29, 0x3B), None),
            ("Dato reale POC gg 15-75", Pt(8.5), RGBColor(0x1E, 0x29, 0x3B), None),
        ]
        add_table_row(table, cells2, "FFFFFF")
        break

print("[OK] Slide 46 (was 45) A4 table updated")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 55 (was 54) — Closing: v4.2 → v4.3, date
# ═══════════════════════════════════════════════════════════════════
slide55 = prs.slides[54]  # was index 53, now 54
for shape in slide55.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'v4.2' in run.text or '4.2' in run.text:
                    run.text = run.text.replace('4.2', '4.3')
                if '26 marzo' in run.text:
                    run.text = run.text.replace('26 marzo', '27 marzo')

print("[OK] Slide 55 (was 54) updated (v4.3, 27 marzo)")


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════
prs.save(PPTX_PATH)
print(f"\n[OK] Saved to {PPTX_PATH}")
print(f"  Total slides: {len(prs.slides)}")
